#!/usr/bin/env python3
"""Generate MedScope AI TFM presentation (Fundae / video defense)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "docs" / "Thesis" / "slides" / "MedScope-AI-TFM.pptx"
SCREENSHOTS = ROOT / "docs" / "figures" / "screenshots"
EDA = ROOT / "docs" / "figures" / "eda"
ICON = ROOT / "frontend" / "public" / "app-icon.png"

PRIMARY = RGBColor(0x00, 0x58, 0xBC)
TEXT = RGBColor(0x19, 0x1C, 0x1D)
MUTED = RGBColor(0x5C, 0x67, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)


def _set_slide_bg(slide, color: RGBColor = WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(slide, prs: Presentation) -> None:
    bar = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.12),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def _add_title(slide, text: str, top: float = 0.45, size: int = 32) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Segoe UI"


def _add_bullets(
    slide,
    items: list[str],
    left: float = 0.75,
    top: float = 1.55,
    width: float = 11.5,
    height: float = 5.0,
    font_size: int = 20,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.font.name = "Segoe UI"
        p.space_after = Pt(10)


def _add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED
    p.font.name = "Segoe UI"


def _add_image_if_exists(slide, path: Path, left, top, width, height=None) -> bool:
    if not path.is_file():
        return False
    if height is None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        slide.shapes.add_picture(
            str(path),
            Inches(left),
            Inches(top),
            width=Inches(width),
            height=Inches(height),
        )
    return True


def _content_slide(prs: Presentation, title: str) -> object:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide)
    _add_header_bar(slide, prs)
    _add_title(slide, title)
    return slide


def _architecture_boxes(slide) -> None:
    boxes = [
        ("Navegador", 0.6),
        ("Vercel\nReact SPA", 3.0),
        ("Render\nFastAPI + ML", 5.8),
        ("Supabase\nPostgreSQL", 8.6),
    ]
    for label, left in boxes:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(2.0), Inches(2.0), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.color.rgb = PRIMARY
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT
        p.font.name = "Segoe UI"

    arrows = [(2.55, 2.45), (5.05, 2.45), (7.85, 2.45)]
    for left, top in arrows:
        arr = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(0.4), Inches(0.3))
        p = arr.text_frame.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(24)
        p.font.color.rgb = PRIMARY


def _ml_pipeline_boxes(slide) -> None:
    steps = [
        "UCI Dataset",
        "Preprocess",
        "Train LR/RF",
        "Evaluate",
        "Serialize",
        "Infer API",
    ]
    left = 0.45
    for step in steps:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(2.0), Inches(1.85), Inches(0.75))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.color.rgb = PRIMARY
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT
        left += 2.05


def _metrics_table(slide) -> None:
    rows, cols = 4, 3
    table = slide.shapes.add_table(rows, cols, Inches(0.9), Inches(1.7), Inches(8.5), Inches(2.2)).table

    headers = ["Métrica", "Logistic Regression", "Random Forest"]
    data = [
        ["Accuracy", "60,7 %", "82,2 %"],
        ["Recall", "54,2 %", "20,1 %"],
        ["ROC-AUC", "0,61", "0,59"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.font.name = "Segoe UI"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = TEXT
                p.font.name = "Segoe UI"
            if c == 0:
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header_bar(slide, prs)
    if ICON.is_file():
        _add_image_if_exists(slide, ICON, 5.9, 0.55, 1.5, 1.5)

    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.8), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "MedScope AI"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Segoe UI"

    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.8), Inches(0.8))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "Sistema de apoyo a la decisión clínica con IA explicable"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT
    p.font.name = "Segoe UI"

    meta = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.8), Inches(1.2))
    tf = meta.text_frame
    for i, line in enumerate(
        [
            "[Tu nombre completo]",
            "Máster Desarrollo con IA — BIG School",
            "Julio 2026",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = MUTED
        p.font.name = "Segoe UI"

    # Slide 2 — Problem
    slide = _content_slide(prs, "El problema clínico")
    _add_bullets(
        slide,
        [
            "Las readmisiones hospitalarias a 30 días son costosas y evitables en parte",
            "Los sistemas HIS almacenan datos pero no siempre anticipan el riesgo",
            "Los modelos de IA sin explicación generan desconfianza en el personal clínico",
            "Necesidad: predicción + transparencia + exploración de escenarios",
        ],
    )

    # Slide 3 — Solution
    slide = _content_slide(prs, "MedScope AI — ¿Qué es?")
    _add_bullets(
        slide,
        [
            "Plataforma web CDSS (Clinical Decision Support System)",
            "Predicción de riesgo de readmisión a 30 días (diabetes, dataset UCI)",
            "SHAP: explicación de qué variables impulsan el riesgo",
            "Simulación what-if: mejor glucosa o menos admisiones previas",
            "No es un sistema de diagnóstico ni de prescripción",
        ],
        width=6.2,
    )
    _add_image_if_exists(slide, SCREENSHOTS / "01_splash.png", 7.1, 1.5, 5.5)

    # Slide 4 — Objectives
    slide = _content_slide(prs, "Objetivos del TFM")
    _add_bullets(
        slide,
        [
            "Pipeline ML reproducible offline → producción",
            "API REST segura (JWT, roles clínico / analista / admin)",
            "Interfaz clínica moderna (React + TypeScript)",
            "Persistencia PostgreSQL e historial auditable",
            "Despliegue cloud real (Vercel + Render + Supabase)",
        ],
    )

    # Slide 5 — Architecture
    slide = _content_slide(prs, "Arquitectura — monolito modular")
    _architecture_boxes(slide)
    _add_bullets(
        slide,
        [
            "Frontend: React en Vercel",
            "Backend + ML: FastAPI en Render (Docker)",
            "Base de datos: PostgreSQL (Supabase)",
            "ML: entrenamiento offline; inferencia al arranque del API",
        ],
        top=3.5,
        font_size=18,
    )
    _add_footer(slide, "Un repositorio, tres capas desacopladas — sin microservicios")

    # Slide 6 — ML Pipeline
    slide = _content_slide(prs, "Pipeline ML — train → serialize → infer")
    _ml_pipeline_boxes(slide)
    _add_bullets(
        slide,
        [
            "Dataset: UCI Diabetes 130-US hospitals (101.766 encuentros)",
            "19 features clínicas · split 80/20 · random_state=42",
            "Modelos: Logistic Regression, Random Forest, XGBoost",
            "Artefactos: model.pkl, preprocessor.pkl, SHAP background",
            "Nunca se reentrena en cada petición",
        ],
        top=3.2,
        font_size=17,
    )

    # Slide 7 — ML Results
    slide = _content_slide(prs, "Resultados ML (test hold-out)")
    _metrics_table(slide)
    _add_bullets(
        slide,
        [
            "Modelo en producción: Logistic Regression v1.0.0",
            "Criterio: priorizar recall frente a accuracy pura",
            "KPI accuracy > 75 % no alcanzado — limitación documentada",
        ],
        left=9.6,
        top=1.7,
        width=3.0,
        height=3.5,
        font_size=14,
    )
    _add_image_if_exists(slide, EDA / "02_target_distribution.png", 9.4, 4.0, 3.3)

    # Slide 8 — Demo transition
    slide = _content_slide(prs, "Demostración en producción")
    demo = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.5))
    p = demo.text_frame.paragraphs[0]
    p.text = "Demo en vivo"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    url = slide.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.6))
    p = url.text_frame.paragraphs[0]
    p.text = "https://medscope-ai-delta.vercel.app"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT

    _add_bullets(
        slide,
        [
            "API: medscope-ai-q8tg.onrender.com",
            "Casos sintéticos · sin datos reales de pacientes",
            "A continuación: recorrido clínico completo en vídeo",
        ],
        top=4.2,
        font_size=18,
    )

    # Slide 9 — Screenshots grid
    slide = _content_slide(prs, "Interfaz clínica — vistas clave")
    grid = [
        ("03_dashboard.png", 0.6, 1.5),
        ("05_prediction_result_shap.png", 6.8, 1.5),
        ("06_simulation.png", 0.6, 4.0),
        ("08_analytics.png", 6.8, 4.0),
    ]
    for name, left, top in grid:
        _add_image_if_exists(slide, SCREENSHOTS / name, left, top, 5.9, 2.3)
    _add_footer(slide, "Capturas reales de la aplicación (Playwright, T-808)")

    # Slide 10 — Quality
    slide = _content_slide(prs, "Calidad de software")
    _add_bullets(
        slide,
        [
            "215+ tests backend (pytest) · cobertura ~95 %",
            "Tests ML (RTS-010) · E2E Playwright flujo MVP completo",
            "CI en GitHub Actions · despliegue automático en main",
            "Esquema PostgreSQL en 3NF · migraciones Alembic",
            "RBAC: admin, clinician, analyst, nurse",
        ],
    )

    # Slide 11 — Conclusions
    slide = _content_slide(prs, "Conclusiones")
    _add_bullets(
        slide,
        [
            "CDSS funcional end-to-end desplegado en cloud",
            "Explicabilidad SHAP y simulación como diferencial frente a un clasificador opaco",
            "Arquitectura enterprise lista para evolución (FHIR, multi-centro)",
            "Transparencia sobre limitaciones del modelo (precision, generalización)",
            "Contribución TFM: ingeniería + IA explicable, no solo accuracy",
        ],
    )

    # Slide 12 — Future work
    slide = _content_slide(prs, "Trabajo futuro")
    _add_bullets(
        slide,
        [
            "Validación externa en otro hospital",
            "Calibración de umbral y nuevas features",
            "Integración FHIR / estándares clínicos",
            "Mejora continua del modelo (ensemble, más datos)",
        ],
        height=3.5,
    )
    thanks = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.8), Inches(0.8))
    p = thanks.text_frame.paragraphs[0]
    p.text = "Gracias por su atención"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    contact = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.8), Inches(0.5))
    p = contact.text_frame.paragraphs[0]
    p.text = "[Tu email / LinkedIn]"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.color.rgb = MUTED

    return prs


def main() -> None:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(OUT_PATH))
    print(f"Generated: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
